#!/usr/bin/env python3
"""Mechanicznie ekstrahuje strukturę prezentacji PPTX bez OCR ani AI.

Wyniki są przeznaczone wyłącznie do lokalnego katalogu ``imports_working/``.
Skrypt zachowuje pochodzenie każdej informacji: prezentacja → slajd → element.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import posixpath
import re
import shutil
import sys
import zipfile
from collections import Counter, defaultdict
from datetime import UTC, datetime
from pathlib import Path
from typing import Any, Iterable
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


PROJECT_ROOT = Path(__file__).resolve().parents[1]
WORKING_ROOT = PROJECT_ROOT / "imports_working"
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}
REL_ID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
REL_EMBED = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"


def sha256_bytes(payload: bytes) -> str:
    return hashlib.sha256(payload).hexdigest()


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def local_name(element: Any) -> str:
    tag = getattr(element, "tag", "")
    return tag.rsplit("}", 1)[-1] if isinstance(tag, str) else "unknown"


def emu_position(shape: Any) -> dict[str, Any]:
    values = {"left": shape.left, "top": shape.top, "width": shape.width, "height": shape.height}
    return {
        "emu": values,
        "inches": {key: round(value / 914400, 4) for key, value in values.items()},
    }


def enum_name(value: Any) -> str:
    return getattr(value, "name", str(value))


def clean_text(value: str | None) -> str:
    return (value or "").replace("\r", "").strip()


def text_from_xml(element: ET.Element) -> str:
    return "".join(node.text or "" for node in element.findall(".//a:t", NS)).strip()


def resolve_part(source_part: str, target: str) -> str:
    return posixpath.normpath(posixpath.join(posixpath.dirname(source_part), target))


def notes_for_slide(archive: zipfile.ZipFile, slide_number: int) -> dict[str, Any]:
    """Odczytuje wyłącznie obszar notatek prezentera z XML, jeśli istnieje."""
    slide_part = f"ppt/slides/slide{slide_number}.xml"
    rels_part = f"ppt/slides/_rels/slide{slide_number}.xml.rels"
    if rels_part not in archive.namelist():
        return {"available": False, "text": "", "part": None}

    rels_root = ET.fromstring(archive.read(rels_part))
    notes_target = None
    for relationship in rels_root.findall("rel:Relationship", NS):
        if relationship.attrib.get("Type", "").endswith("/notesSlide"):
            notes_target = relationship.attrib.get("Target")
            break
    if not notes_target:
        return {"available": False, "text": "", "part": None}

    notes_part = resolve_part(slide_part, notes_target)
    if notes_part not in archive.namelist():
        return {"available": False, "text": "", "part": notes_part}

    root = ET.fromstring(archive.read(notes_part))
    texts: list[str] = []
    for shape in root.findall(".//p:sp", NS):
        placeholder = shape.find(".//p:nvPr/p:ph", NS)
        if placeholder is not None and placeholder.attrib.get("type", "body") == "body":
            candidate = text_from_xml(shape)
            if candidate:
                texts.append(candidate)
    return {"available": True, "text": "\n".join(texts), "part": notes_part}


def hyperlink_relationships_for_slide(archive: zipfile.ZipFile, slide_number: int) -> list[dict[str, str]]:
    """Zwraca relacje hyperlinków, także gdy biblioteka nie mapuje ich na kształt."""
    rels_part = f"ppt/slides/_rels/slide{slide_number}.xml.rels"
    if rels_part not in archive.namelist():
        return []
    root = ET.fromstring(archive.read(rels_part))
    return [
        {
            "relationship_id": relationship.attrib["Id"],
            "target": relationship.attrib.get("Target", ""),
            "target_mode": relationship.attrib.get("TargetMode", ""),
        }
        for relationship in root.findall("rel:Relationship", NS)
        if relationship.attrib.get("Type", "").endswith("/hyperlink")
    ]


def shape_relationship_links(shape: Any) -> list[dict[str, str]]:
    """Zbiera kliknięcia/hyperlinki zapisane w XML elementu."""
    results: list[dict[str, str]] = []
    for node in shape._element.iter():
        node_type = local_name(node)
        if node_type not in {"hlinkClick", "hlinkHover"}:
            continue
        relation_id = node.attrib.get(REL_ID)
        if not relation_id:
            continue
        try:
            relation = shape.part.rels[relation_id]
            target = relation.target_ref if relation.is_external else str(relation.target_part.partname)
        except (KeyError, AttributeError, ValueError):
            target = f"nierozwiązana_relacja:{relation_id}"
        results.append({"kind": node_type, "target": target, "relationship_id": relation_id})
    return results


def run_hyperlinks(shape: Any) -> list[dict[str, Any]]:
    links: list[dict[str, Any]] = []
    if not getattr(shape, "has_text_frame", False):
        return links
    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
        for run_index, run in enumerate(paragraph.runs, start=1):
            address = run.hyperlink.address
            if address:
                links.append(
                    {
                        "kind": "run_hyperlink",
                        "target": address,
                        "text": run.text,
                        "paragraph_index": paragraph_index,
                        "run_index": run_index,
                    }
                )
    return links


def text_content(shape: Any) -> dict[str, Any] | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    paragraphs: list[dict[str, Any]] = []
    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
        runs = []
        for run_index, run in enumerate(paragraph.runs, start=1):
            runs.append(
                {
                    "run_index": run_index,
                    "text": run.text,
                    "hyperlink": run.hyperlink.address,
                }
            )
        paragraphs.append(
            {
                "paragraph_index": paragraph_index,
                "level": paragraph.level,
                "text": "".join(run["text"] for run in runs),
                "runs": runs,
            }
        )
    text = "\n".join(paragraph["text"] for paragraph in paragraphs).strip()
    return {"text": text, "paragraphs": paragraphs}


def table_content(shape: Any) -> dict[str, Any] | None:
    if not getattr(shape, "has_table", False):
        return None
    table = shape.table
    rows = []
    for row_index, row in enumerate(table.rows, start=1):
        cells = []
        for column_index, cell in enumerate(row.cells, start=1):
            cells.append({"row": row_index, "column": column_index, "text": cell.text})
        rows.append(cells)
    return {"rows": len(table.rows), "columns": len(table.columns), "cells": rows}


def chart_content(shape: Any) -> dict[str, Any] | None:
    if not getattr(shape, "has_chart", False):
        return None
    chart = shape.chart
    result: dict[str, Any] = {"chart_type": enum_name(chart.chart_type)}
    try:
        result["series_count"] = len(chart.series)
    except (AttributeError, TypeError):
        result["series_count"] = None
    return result


def embedded_image_relations(shape: Any) -> list[dict[str, Any]]:
    """Zwraca osadzone obrazy z XML, również w placeholderach typu Picture."""
    result = []
    for node in shape._element.iter():
        if local_name(node) != "blip":
            continue
        relationship_id = node.attrib.get(REL_EMBED) or node.attrib.get(REL_ID)
        if not relationship_id:
            continue
        try:
            relationship = shape.part.rels[relationship_id]
            target_part = relationship.target_part
            result.append(
                {
                    "relationship_id": relationship_id,
                    "source_part": str(target_part.partname),
                    "content_type": getattr(target_part, "content_type", None),
                    "payload": target_part.blob,
                }
            )
        except (AttributeError, KeyError, ValueError):
            result.append({"relationship_id": relationship_id, "extraction_error": "nierozwiązana_relacja_obrazu"})
    return result


def image_contents(shape: Any, images_directory: Path, slide_number: int, image_counter: list[int]) -> list[dict[str, Any]]:
    """Wyodrębnia każdy obraz przypisany do kształtu i zachowuje jego relację XML."""
    records = []
    relations = embedded_image_relations(shape)
    for relation in relations:
        if "extraction_error" in relation:
            records.append(relation)
            continue
        source_part = relation["source_part"]
        extension = re.sub(r"[^a-zA-Z0-9]", "", Path(source_part).suffix.lower()) or "bin"
        image_counter[0] += 1
        filename = f"slide_{slide_number:03d}_image_{image_counter[0]:02d}.{extension}"
        payload = relation.pop("payload")
        (images_directory / filename).write_bytes(payload)
        records.append(
            {
                "filename": filename,
                "extension": extension,
                "sha256": sha256_bytes(payload),
                "bytes": len(payload),
                **relation,
            }
        )
    return records


def media_relationships(shape: Any) -> list[dict[str, str]]:
    media: list[dict[str, str]] = []
    relationship_ids = {node.attrib[REL_ID] for node in shape._element.iter() if REL_ID in node.attrib}
    for relationship_id in relationship_ids:
        try:
            relation = shape.part.rels[relationship_id]
        except KeyError:
            continue
        relation_type = getattr(relation, "reltype", "")
        if not any(token in relation_type for token in ("/audio", "/video", "/media")):
            continue
        try:
            target = relation.target_ref if relation.is_external else str(relation.target_part.partname)
        except (AttributeError, ValueError):
            target = "nierozpoznany_cel"
        media.append(
            {"relationship_id": relationship_id, "relationship_type": relation_type, "target": target}
        )
    return media


def extract_shape(
    shape: Any,
    slide_number: int,
    images_directory: Path,
    image_counter: list[int],
    group_path: list[int] | None = None,
) -> list[dict[str, Any]]:
    """Ekstrahuje element; dla grup zachowuje element nadrzędny i jego dzieci."""
    group_path = group_path or []
    shape_type = enum_name(shape.shape_type)
    element: dict[str, Any] = {
        "element_id": f"slide_{slide_number:03d}_shape_{shape.shape_id}",
        "shape_id": shape.shape_id,
        "name": shape.name,
        "shape_type": shape_type,
        "xml_type": local_name(shape._element),
        "group_path": group_path,
        "position": emu_position(shape),
        "text": text_content(shape),
        "hyperlinks": run_hyperlinks(shape) + shape_relationship_links(shape),
        "table": table_content(shape),
        "chart": chart_content(shape),
        "media_relationships": media_relationships(shape),
    }
    element["images"] = image_contents(shape, images_directory, slide_number, image_counter)

    elements = [element]
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        element["group_children"] = []
        for child in shape.shapes:
            child_elements = extract_shape(
                child, slide_number, images_directory, image_counter, group_path + [shape.shape_id]
            )
            element["group_children"].extend(child["element_id"] for child in child_elements)
            elements.extend(child_elements)
    return elements


def xml_features(archive: zipfile.ZipFile, slide_number: int) -> dict[str, Any]:
    part = f"ppt/slides/slide{slide_number}.xml"
    root = ET.fromstring(archive.read(part))
    tags = Counter(local_name(node) for node in root.iter())
    namespaces = [node.tag.rsplit("}", 1)[0].lstrip("{") for node in root.iter() if isinstance(node.tag, str)]
    return {
        "xml_element_counts": dict(sorted(tags.items())),
        # Deklaracja przestrzeni nazw diagramu jest obecna w wielu zwykłych
        # slajdach, dlatego wykrywamy rzeczywisty element diagramu, nie tekst
        # w nagłówku XML.
        "has_smartart_or_diagram": any("drawingml/2006/diagram" in namespace for namespace in namespaces),
        "has_ole_object": any(local_name(node) == "oleObj" for node in root.iter()),
        "has_embedded_media": any(local_name(node) in {"videoFile", "audioFile"} for node in root.iter()),
    }


def choose_title(elements: Iterable[dict[str, Any]]) -> str | None:
    candidates = []
    for element in elements:
        text = (element.get("text") or {}).get("text", "")
        if not text:
            continue
        name = (element.get("name") or "").lower()
        shape_type = (element.get("shape_type") or "").lower()
        priority = 0 if "title" in name or "placeholder" in shape_type else 1
        candidates.append((priority, element["position"]["emu"]["top"], element["position"]["emu"]["left"], text))
    return min(candidates)[-1] if candidates else None


def aggregate_hyperlinks(elements: Iterable[dict[str, Any]]) -> list[dict[str, Any]]:
    """Łączy podwójne reprezentacje tego samego linku (run i XML)."""
    grouped: dict[tuple[str, str], dict[str, Any]] = {}
    for element in elements:
        element_links = element.get("hyperlinks", [])
        xml_targets = {link.get("target") for link in element_links if link.get("kind") in {"hlinkClick", "hlinkHover"}}
        for link in element_links:
            target = link.get("target")
            if not target:
                continue
            # Link w runie i hlinkClick w XML opisują zwykle tę samą relację.
            # Do licznika używamy relacji XML, a run pozostaje szczegółem pola tekstowego.
            if link.get("kind") == "run_hyperlink" and target in xml_targets:
                continue
            relation_or_run = link.get("relationship_id") or (
                f"run:{link.get('paragraph_index')}:{link.get('run_index')}"
            )
            key = (element["element_id"], relation_or_run)
            if key not in grouped:
                grouped[key] = {
                    "element_id": element["element_id"],
                    "target": target,
                    "sources": [],
                }
            grouped[key]["sources"].append({key: value for key, value in link.items() if key != "target"})
    return list(grouped.values())


def include_unmapped_hyperlinks(
    links: list[dict[str, Any]], relationship_inventory: Iterable[dict[str, str]]
) -> list[dict[str, Any]]:
    mapped_ids = {
        source.get("relationship_id")
        for link in links
        for source in link.get("sources", [])
        if source.get("relationship_id")
    }
    for relation in relationship_inventory:
        if relation["relationship_id"] in mapped_ids:
            continue
        links.append(
            {
                "element_id": None,
                "target": relation["target"],
                "sources": [{"kind": "slide_relationship_not_mapped_to_shape", **relation}],
            }
        )
    return links


def slide_diagnostics(slide: dict[str, Any], slide_width: int, slide_height: int) -> dict[str, Any]:
    elements = slide["elements"]
    text_elements = [item for item in elements if (item.get("text") or {}).get("text")]
    image_records = [
        image
        for item in elements
        for image in item.get("images", [])
        if not image.get("extraction_error")
    ]
    image_elements = [item for item in elements if any(not image.get("extraction_error") for image in item.get("images", []))]
    tables = [item for item in elements if item.get("table")]
    groups = [item for item in elements if item.get("shape_type") == "GROUP"]
    text_characters = sum(len((item.get("text") or {}).get("text", "")) for item in text_elements)
    image_area = sum(item["position"]["emu"]["width"] * item["position"]["emu"]["height"] for item in image_elements)
    image_coverage = round(image_area / (slide_width * slide_height), 3) if slide_width and slide_height else 0
    reasons: list[str] = []
    if not text_elements and image_records:
        reasons.append("brak_tekstu_i_obraz")
    if image_records and text_characters <= 40:
        reasons.append("obraz_z_niewielka_iloscia_tekstu")
    if image_coverage >= 0.45:
        reasons.append("duzy_udzial_powierzchni_obrazow")
    if slide["xml_features"]["has_smartart_or_diagram"]:
        reasons.append("smartart_lub_diagram")
    if any(item.get("chart") for item in elements):
        reasons.append("wykres")
    if groups:
        reasons.append("zgrupowane_elementy")
    if len(elements) >= 12:
        reasons.append("wiele_elementow")
    return {
        "text_characters": text_characters,
        "text_element_count": len(text_elements),
        "image_count": len(image_records),
        "table_count": len(tables),
        "group_count": len(groups),
        "image_coverage_estimate": image_coverage,
        "requires_visual_review": bool(reasons),
        "visual_review_reasons": reasons,
        "possible_reading_order_issue": len(text_elements) > 1,
    }


def choose_validation_sample(slides: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Wybiera do ośmiu różnych slajdów o cechach wymaganych do kontroli ręcznej."""
    selections: list[tuple[str, dict[str, Any]]] = []
    seen: set[int] = set()

    def add(label: str, candidates: Iterable[dict[str, Any]]) -> None:
        for slide in candidates:
            if slide["slide_number"] not in seen:
                seen.add(slide["slide_number"])
                selections.append((label, slide))
                return

    if slides:
        add("pierwszy", [slides[0]])
        add("ostatni", [slides[-1]])

    add(
        "typowy_tekstowy",
        sorted(
            (s for s in slides if s["diagnostics"]["text_characters"] > 0 and s["diagnostics"]["image_count"] == 0),
            key=lambda s: (abs(s["diagnostics"]["text_characters"] - 250), s["slide_number"]),
        ),
    )
    add(
        "duzo_tekstu",
        sorted(slides, key=lambda s: (-s["diagnostics"]["text_characters"], s["slide_number"])),
    )
    add(
        "jeden_duzy_obraz",
        sorted(
            (s for s in slides if s["diagnostics"]["image_count"] == 1),
            key=lambda s: (-s["diagnostics"]["image_coverage_estimate"], s["slide_number"]),
        ),
    )
    add(
        "kilka_obrazow",
        sorted(
            (s for s in slides if s["diagnostics"]["image_count"] >= 2),
            key=lambda s: (-s["diagnostics"]["image_count"], s["slide_number"]),
        ),
    )
    add(
        "najbardziej_zlozony",
        sorted(
            slides,
            key=lambda s: (
                -(len(s["elements"]) + 3 * s["diagnostics"]["image_count"] + 3 * s["diagnostics"]["table_count"]),
                s["slide_number"],
            ),
        ),
    )
    add(
        "nietypowy",
        sorted(
            (s for s in slides if s["diagnostics"]["table_count"] or s["notes"]["text"] or s["diagnostics"]["requires_visual_review"]),
            key=lambda s: (
                -s["diagnostics"]["table_count"],
                -bool(s["notes"]["text"]),
                -s["diagnostics"]["image_coverage_estimate"],
                s["slide_number"],
            ),
        ),
    )
    add("dodatkowy_reprezentatywny", slides)
    return [
        {"purpose": label, "slide_number": slide["slide_number"], "title": slide["title"]}
        for label, slide in selections[:8]
    ]


def markdown_for_slides(document: dict[str, Any]) -> str:
    lines = [f"# {document['presentation']['filename']}", ""]
    for slide in document["slides"]:
        lines.extend([f"# Slajd {slide['slide_number']:03d}", ""])
        lines.extend(["## Tytuł", slide["title"] or "_Nie wykryto tytułu._", ""])
        lines.append("## Tekst")
        text_elements = [element for element in slide["elements"] if (element.get("text") or {}).get("text")]
        if text_elements:
            for element in sorted(
                text_elements,
                key=lambda item: (item["position"]["emu"]["top"], item["position"]["emu"]["left"], item["shape_id"]),
            ):
                lines.append(f"- `{element['element_id']}`: {(element['text']['text'])}")
        else:
            lines.append("_Brak tekstu wykrywalnego przez bibliotekę._")
        lines.extend(["", "## Tabele"])
        table_elements = [element for element in slide["elements"] if element.get("table")]
        if table_elements:
            for element in table_elements:
                lines.append(f"### {element['element_id']}")
                for row in element["table"]["cells"]:
                    lines.append(" | ".join(cell["text"].replace("\n", " ") for cell in row))
        else:
            lines.append("_Brak tabel._")
        lines.extend(["", "## Hiperłącza"])
        links = slide["hyperlinks"]
        if links:
            for link in links:
                lines.append(f"- `{link['element_id']}` → {link['target']}")
        else:
            lines.append("_Brak wykrytych hiperłączy._")
        lines.extend(["", "## Notatki", slide["notes"]["text"] or "_Brak wykrywalnych notatek prezentera._", ""])
        lines.append("## Obrazy")
        if slide["images"]:
            for image in slide["images"]:
                lines.append(f"- `{image['element_id']}` → `images/{image['filename']}`")
        else:
            lines.append("_Brak obrazów wykrywalnych przez bibliotekę._")
        lines.extend(["", "## Diagnostyka"])
        diagnostics = slide["diagnostics"]
        if diagnostics["visual_review_reasons"]:
            lines.append("- Wymaga analizy wizualnej: " + ", ".join(diagnostics["visual_review_reasons"]))
        if diagnostics["possible_reading_order_issue"]:
            lines.append("- Kolejność tekstu została ustalona heurystycznie według pozycji pionowej i poziomej.")
        lines.append("")
    return "\n".join(lines)


def diagnostic_report(document: dict[str, Any]) -> str:
    summary = document["diagnostic_summary"]
    lines = [
        "# Raport diagnostyczny PPTX",
        "",
        "## Prezentacja",
        f"- Nazwa: `{document['presentation']['filename']}`",
        f"- Rozmiar: {document['presentation']['bytes']:,} B",
        f"- SHA-256: `{document['presentation']['sha256']}`",
        f"- Slajdy: {summary['slide_count']}",
        f"- Slajdy z tekstem w polach tekstowych: {summary['slides_with_text']}",
        f"- Slajdy z treścią tekstową po uwzględnieniu tabel: {summary['slides_with_text_or_table']}",
        f"- Slajdy z obrazami: {summary['slides_with_images']}",
        f"- Tabele: {summary['table_count']}",
        f"- Hiperłącza: {summary['hyperlink_count']}",
        f"- Slajdy z notatkami: {summary['slides_with_notes']}",
        f"- Wyodrębnione wystąpienia obrazów: {summary['extracted_image_occurrences']}",
        "",
        "## Co odzyskano mechanicznie",
        "- Numery slajdów, układy, elementy, pozycje, pola tekstowe, tabele, wykrywalne hiperłącza, notatki prezentera i obrazy osadzone jako kształty typu Picture.",
        "- Tekst jest zapisany osobno dla każdego elementu i zachowuje identyfikator elementu; w Markdown jest dodatkowo sortowany heurystycznie według położenia.",
        "",
        "## Ograniczenia i ryzyka",
        "- Kolejność tekstu w pliku PowerPoint jest przede wszystkim kolejnością warstw, a nie gwarantowaną kolejnością czytania. Heurystyka pozycji może być błędna w układach wielokolumnowych, nakładających się lub grupowanych.",
        "- SmartArt, diagramy, wykresy, obiekty OLE i multimedia są oznaczane diagnostycznie; ich znaczenie nie jest rekonstruowane jako pełna treść semantyczna.",
        "- Tekst zawarty w obrazach, skanach, RTG/MRI, zrzutach ekranu, strzałkach i schematach nie jest odczytywany, ponieważ importer celowo nie wykonuje OCR ani interpretacji obrazu.",
        "- Automatycznie wyodrębnione obrazy nie są ocenione pod kątem źródła, licencji ani dopuszczalności publikacji.",
        "",
        "## Układy slajdów",
    ]
    for layout, count in summary["layouts"].items():
        lines.append(f"- `{layout}`: {count}")
    lines.extend(["", "## Elementy potencjalnie trudne do ekstrakcji"])
    for feature, count in summary["feature_counts"].items():
        lines.append(f"- {feature}: {count}")
    lines.extend(["", "## Potencjalne duplikaty obrazów"])
    if summary["duplicate_images"]:
        for item in summary["duplicate_images"]:
            lines.append(f"- SHA-256 `{item['sha256']}`: {', '.join(item['occurrences'])}")
    else:
        lines.append("- Nie wykryto powtórzeń wśród obrazów wyodrębnionych jako elementy Picture.")
    lines.extend(["", "## Slajdy wymagające analizy wizualnej"])
    if summary["visual_review_slides"]:
        for item in summary["visual_review_slides"]:
            lines.append(f"- Slajd {item['slide_number']:03d}: {', '.join(item['reasons'])}")
    else:
        lines.append("- Heurystyka nie wykryła oczywistych sygnałów; nie jest to dowód, że analiza wizualna jest zbędna.")
    lines.extend(["", "## Zalecana próbka do ręcznej walidacji"])
    for item in document["validation_sample"]:
        lines.append(f"- Slajd {item['slide_number']:03d} ({item['purpose']}): {item['title'] or 'bez wykrytego tytułu'}")
    lines.extend(["", "## Wyniki niemożliwe do odczytania lub niejednoznaczne"])
    if summary["extraction_errors"]:
        lines.extend(f"- {error}" for error in summary["extraction_errors"])
    else:
        lines.append("- Brak błędów krytycznych biblioteki podczas tej ekstrakcji.")
    return "\n".join(lines) + "\n"


def extract(pptx_path: Path, output_directory: Path, overwrite: bool) -> dict[str, Any]:
    pptx_path = pptx_path.resolve()
    if not pptx_path.is_file() or pptx_path.suffix.lower() != ".pptx":
        raise ValueError(f"Nieprawidłowy plik PPTX: {pptx_path}")
    output_directory = output_directory.resolve()
    try:
        output_directory.relative_to(WORKING_ROOT.resolve())
    except ValueError as error:
        raise ValueError(f"Wyniki mogą być zapisane wyłącznie w {WORKING_ROOT}") from error
    if output_directory.exists() and any(output_directory.iterdir()):
        if not overwrite:
            raise ValueError(f"Katalog wynikowy nie jest pusty: {output_directory}. Użyj --overwrite.")
        shutil.rmtree(output_directory)
    images_directory = output_directory / "images"
    images_directory.mkdir(parents=True, exist_ok=True)

    presentation = Presentation(pptx_path)
    extracted_at = datetime.now(UTC).replace(microsecond=0).isoformat().replace("+00:00", "Z")
    document: dict[str, Any] = {
        "schema_version": "1.0",
        "extracted_at": extracted_at,
        "presentation": {
            "filename": pptx_path.name,
            "bytes": pptx_path.stat().st_size,
            "sha256": sha256_file(pptx_path),
            "slide_width_emu": presentation.slide_width,
            "slide_height_emu": presentation.slide_height,
        },
        "slides": [],
        "package_media": [],
    }
    extraction_errors: list[str] = []
    with zipfile.ZipFile(pptx_path) as archive:
        document["package_media"] = [
            {"path": name, "bytes": archive.getinfo(name).file_size, "sha256": sha256_bytes(archive.read(name))}
            for name in archive.namelist()
            if name.startswith("ppt/media/")
        ]
        for slide_number, slide in enumerate(presentation.slides, start=1):
            try:
                image_counter = [0]
                elements: list[dict[str, Any]] = []
                for shape in slide.shapes:
                    elements.extend(extract_shape(shape, slide_number, images_directory, image_counter))
                text_elements = [element for element in elements if (element.get("text") or {}).get("text")]
                for reading_order, element in enumerate(
                    sorted(
                        text_elements,
                        key=lambda item: (
                            item["position"]["emu"]["top"],
                            item["position"]["emu"]["left"],
                            item["shape_id"],
                        ),
                    ),
                    start=1,
                ):
                    element["reading_order_by_position"] = reading_order
                notes = notes_for_slide(archive, slide_number)
                slide_record: dict[str, Any] = {
                    "presentation_filename": pptx_path.name,
                    "slide_number": slide_number,
                    "slide_layout": slide.slide_layout.name,
                    "title": choose_title(elements),
                    "elements": elements,
                    "notes": notes,
                    "xml_features": xml_features(archive, slide_number),
                }
                slide_record["hyperlinks"] = include_unmapped_hyperlinks(
                    aggregate_hyperlinks(elements), hyperlink_relationships_for_slide(archive, slide_number)
                )
                slide_record["images"] = [
                {"element_id": element["element_id"], **image}
                for element in elements
                for image in element.get("images", [])
                if not image.get("extraction_error")
                ]
                slide_record["diagnostics"] = slide_diagnostics(
                    slide_record, presentation.slide_width, presentation.slide_height
                )
                document["slides"].append(slide_record)
            except Exception as error:  # Nie ukrywamy błędu pojedynczego slajdu przed raportem.
                extraction_errors.append(f"Slajd {slide_number:03d}: {type(error).__name__}: {error}")

    layouts = Counter(slide["slide_layout"] for slide in document["slides"])
    image_occurrences: dict[str, list[str]] = defaultdict(list)
    feature_counts = Counter()
    visual_review_slides = []
    for slide in document["slides"]:
        for image in slide["images"]:
            image_occurrences[image["sha256"]].append(f"slajd_{slide['slide_number']:03d}/{image['filename']}")
        if slide["xml_features"]["has_smartart_or_diagram"]:
            feature_counts["slajdy_z_SmartArt_lub_diagramem"] += 1
        if slide["xml_features"]["has_ole_object"]:
            feature_counts["slajdy_z_obiektem_OLE"] += 1
        if slide["xml_features"]["has_embedded_media"]:
            feature_counts["slajdy_z_osadzonym_multimedium"] += 1
        if any(element.get("chart") for element in slide["elements"]):
            feature_counts["slajdy_z_wykresem"] += 1
        if slide["diagnostics"]["group_count"]:
            feature_counts["slajdy_z_grupami_elementow"] += 1
        if slide["diagnostics"]["requires_visual_review"]:
            visual_review_slides.append(
                {"slide_number": slide["slide_number"], "reasons": slide["diagnostics"]["visual_review_reasons"]}
            )
    document["diagnostic_summary"] = {
        "slide_count": len(document["slides"]),
        "slides_with_text": sum(bool(slide["diagnostics"]["text_characters"]) for slide in document["slides"]),
        "slides_with_text_or_table": sum(
            bool(slide["diagnostics"]["text_characters"] or slide["diagnostics"]["table_count"])
            for slide in document["slides"]
        ),
        "slides_with_images": sum(bool(slide["images"]) for slide in document["slides"]),
        "table_count": sum(slide["diagnostics"]["table_count"] for slide in document["slides"]),
        "hyperlink_count": sum(len(slide["hyperlinks"]) for slide in document["slides"]),
        "slides_with_notes": sum(bool(slide["notes"]["text"]) for slide in document["slides"]),
        "extracted_image_occurrences": sum(len(slide["images"]) for slide in document["slides"]),
        "layouts": dict(sorted(layouts.items())),
        "feature_counts": dict(sorted(feature_counts.items())),
        "duplicate_images": [
            {"sha256": digest, "occurrences": occurrences}
            for digest, occurrences in image_occurrences.items()
            if len(occurrences) > 1
        ],
        "visual_review_slides": visual_review_slides,
        "extraction_errors": extraction_errors,
    }
    document["validation_sample"] = choose_validation_sample(document["slides"])

    (output_directory / "slides.json").write_text(
        json.dumps(document, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    (output_directory / "slides.md").write_text(markdown_for_slides(document), encoding="utf-8")
    (output_directory / "diagnostic_report.md").write_text(diagnostic_report(document), encoding="utf-8")
    return document


def parse_arguments() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", type=Path, help="Lokalny plik .pptx do analizy")
    parser.add_argument(
        "--output",
        type=Path,
        help="Katalog w imports_working/. Domyślnie imports_working/<nazwa-prezentacji>/.",
    )
    parser.add_argument("--overwrite", action="store_true", help="Zastępuje istniejący katalog wynikowy")
    return parser.parse_args()


def main() -> int:
    arguments = parse_arguments()
    output = arguments.output or WORKING_ROOT / arguments.pptx.stem
    try:
        document = extract(arguments.pptx, output, arguments.overwrite)
    except (OSError, ValueError, zipfile.BadZipFile) as error:
        print(f"BŁĄD: {error}", file=sys.stderr)
        return 1
    summary = document["diagnostic_summary"]
    print(f"Przetworzono: {document['presentation']['filename']}")
    print(f"Slajdy: {summary['slide_count']}; obrazy: {summary['extracted_image_occurrences']}; tabele: {summary['table_count']}")
    print(f"Wyniki: {output.resolve()}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
